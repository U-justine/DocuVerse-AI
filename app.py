#!/usr/bin/env python3
"""
🌟 DOCUVERSE AI 🌟
Revolutionary Document Assistant with stunning design
Copyright © 2025 Justine & Krishna. All Rights Reserved.
"""

import os
import re
import time
from datetime import datetime
from typing import List
from difflib import SequenceMatcher
from io import BytesIO
import importlib.util

import requests
import streamlit as st
from markupsafe import escape
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
from transformers import pipeline, AutoTokenizer
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer as SumyTokenizer
from sumy.summarizers.lsa import LsaSummarizer
import nltk
import cohere
from huggingface_hub import login as hf_login
from dotenv import load_dotenv

# Load environment variables and NLTK data
load_dotenv()
nltk.download('punkt', quiet=True)
nltk.download('punkt_tab', quiet=True)

# -----------------------------
# Page Configuration
# -----------------------------
st.set_page_config(
    page_title="DocuVerse AI - Revolutionary Document Assistant",
    page_icon="🌌",
    layout="wide"
)

# Debugging flags
SHOW_KEY_BANNERS = False
DEBUG_API_ERRORS = False  # Suppress error messages in production

# -----------------------------
# Styling
# -----------------------------
def load_revolutionary_css():
    """Load custom CSS for the Streamlit app."""
    css = """
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@300;400;700;900&family=Rajdhani:wght@300;400;600;700&family=Space+Mono:wght@400;700&display=swap');
    .stApp {
        background: linear-gradient(135deg, #0F0C29 0%, #24243e 30%, #302B63 70%, #0F0C29 100%);
        background-attachment: fixed;
        color: #E2E8F0;
        font-family: 'Rajdhani', sans-serif;
    }
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    .stDeployButton {display: none;}
    .main-title {
        font-family: 'Orbitron', monospace;
        font-size: 4.2rem;
        font-weight: 900;
        background: linear-gradient(45deg, #FF6B6B, #4ECDC4, #45B7D1, #96CEB4, #FFEAA7, #FF6B6B);
        background-size: 400% 400%;
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        animation: gradientFlow 4s ease-in-out infinite;
        text-align: center;
        margin: 1.7rem 0 0.4rem 0;
        letter-spacing: 3px;
        text-shadow: 0 0 50px rgba(255, 107, 107, 0.3);
        position: relative;
    }
    .main-title::after {
        content: '';
        position: absolute;
        bottom: -8px;
        left: 50%;
        transform: translateX(-50%);
        width: 180px;
        height: 3px;
        background: linear-gradient(90deg, transparent, #4ECDC4, transparent);
        animation: lineGlow 2s ease-in-out infinite;
    }
    @keyframes gradientFlow {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    @keyframes lineGlow {
        0%, 100% { opacity: 0.3; width: 90px; }
        50% { opacity: 1; width: 260px; }
    }
    .subtitle {
        font-family: 'Rajdhani', sans-serif;
        font-size: 1.4rem;
        font-weight: 300;
        color: #A8A8B3;
        text-align: center;
        margin-bottom: 2rem;
        text-transform: uppercase;
        letter-spacing: 3px;
    }
    .content-section {
        background: linear-gradient(135deg, rgba(255, 255, 255, 0.05) 0%, rgba(255, 255, 255, 0.01) 100%);
        backdrop-filter: blur(15px);
        border-radius: 22px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        margin: 1.2rem 0;
        box-shadow: 0 25px 50px rgba(0, 0, 0, 0.3), inset 0 1px 0 rgba(255, 255, 255, 0.1);
        position: relative;
        overflow: hidden;
    }
    .section-title {
        font-family: 'Orbitron', monospace;
        font-size: 2.2rem;
        font-weight: 700;
        background: linear-gradient(135deg, #4ECDC4, #45B7D1);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        margin-bottom: 1.2rem;
        text-align: center;
        letter-spacing: 2px;
    }
    .cyber-card {
        background: linear-gradient(135deg, rgba(255, 255, 255, 0.08) 0%, rgba(255, 255, 255, 0.02) 100%);
        backdrop-filter: blur(12px);
        border-radius: 16px;
        border: 1px solid rgba(255, 255, 255, 0.15);
        padding: 1rem 1.2rem;
        margin: 0.8rem 0;
        box-shadow: 0 15px 35px rgba(0, 0, 0, 0.2), inset 0 1px 0 rgba(255, 255, 255, 0.1);
        transition: all 0.3s ease;
        position: relative;
        overflow: hidden;
    }
    .cyber-card:hover { transform: translateY(-4px) scale(1.005); }
    .cyber-text {
        background: linear-gradient(135deg, #4ECDC4, #45B7D1);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
    }
    .metric-card {
        background: linear-gradient(135deg, rgba(0, 255, 127, 0.08), rgba(0, 191, 255, 0.08));
        border: 1px solid rgba(0, 255, 127, 0.2);
        border-radius: 16px;
        padding: 1rem;
        text-align: center;
        backdrop-filter: blur(10px);
    }
    .metric-value {
        font-family: 'Orbitron', monospace;
        font-size: 1.8rem;
        font-weight: 700;
        color: #00FF7F;
        text-shadow: 0 0 20px rgba(0, 255, 127, 0.5);
        margin-bottom: 0.3rem;
    }
    .metric-label {
        color: #A8A8B3;
        text-transform: uppercase;
        font-size: 0.85rem;
        font-weight: 600;
        letter-spacing: 1px;
    }
    .keyword-tag {
        background: linear-gradient(135deg, #667eea, #764ba2);
        color: white;
        padding: 0.35rem 0.9rem;
        margin: 0.25rem;
        border-radius: 18px;
        display: inline-block;
        font-size: 0.88rem;
        font-weight: 500;
        box-shadow: 0 5px 15px rgba(102, 126, 234, 0.3);
    }
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%) !important;
        border: none !important;
        border-radius: 14px !important;
        padding: 0.8rem 1.6rem !important;
        color: white !important;
        font-family: 'Rajdhani', sans-serif !important;
        font-size: 1.02rem !important;
        font-weight: 600 !important;
        text-transform: uppercase !important;
        letter-spacing: 1px !important;
        cursor: pointer !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 15px 35px rgba(102, 126, 234, 0.4) !important;
    }
    .stButton > button:hover { transform: translateY(-3px) scale(1.03) !important; }
    .stDownloadButton > button {
        background: linear-gradient(135deg, #4ECDC4 0%, #44A08D 100%) !important;
        border: none !important;
        border-radius: 14px !important;
        color: white !important;
        font-weight: 600 !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 10px 25px rgba(78, 205, 196, 0.3) !important;
    }
    .stDownloadButton > button:hover { transform: translateY(-2px) scale(1.05) !important; }
    .fulltext-container {
        border: 1px solid rgba(255,255,255,0.15);
        background: rgba(255,255,255,0.02);
        border-radius: 14px;
        padding: 1rem;
        height: 520px;
        overflow: auto;
        white-space: pre-wrap;
        line-height: 1.65;
        font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono", "Courier New", monospace;
        font-size: 0.98rem;
    }
    </style>
    """
    st.markdown(css, unsafe_allow_html=True)

# -----------------------------
# Secrets/Keys
# -----------------------------
def get_secret(key: str, default: str = "") -> str:
    """Retrieve secret from environment variables or Streamlit secrets."""
    return os.environ.get(key, st.secrets.get(key, default))

HF_API_KEY = get_secret("HF_API_KEY", "")
COHERE_API_KEY = get_secret("COHERE_API_KEY", "")

# Authenticate Hugging Face silently
if HF_API_KEY:
    try:
        hf_login(token=HF_API_KEY)
    except Exception:
        HF_API_KEY = ""

# -----------------------------
# Cached Resources
# -----------------------------
@st.cache_resource
def get_hf_tokenizer():
    """Load Hugging Face tokenizer."""
    try:
        return AutoTokenizer.from_pretrained("facebook/bart-large-cnn", token=HF_API_KEY if HF_API_KEY else None)
    except Exception:
        return None

@st.cache_resource
def get_st_model():
    """Load SentenceTransformer model."""
    try:
        return SentenceTransformer("all-MiniLM-L6-v2")
    except Exception:
        return None

@st.cache_resource
def get_hf_summarizer():
    """Load Hugging Face summarization pipeline."""
    try:
        return pipeline("summarization", model="facebook/bart-large-cnn", device=-1, token=HF_API_KEY if HF_API_KEY else None)
    except Exception:
        return None

@st.cache_resource
def get_hf_qa():
    """Load Hugging Face question-answering pipeline."""
    try:
        return pipeline("question-answering", model="deepset/roberta-base-squad2", device=-1, token=HF_API_KEY if HF_API_KEY else None)
    except Exception:
        return None

hf_tokenizer = get_hf_tokenizer()

# -----------------------------
# Utilities
# -----------------------------
def split_sentences(text: str) -> List[str]:
    """Split text into sentences."""
    return [s.strip() for s in re.split(r'(?<=[.!?])\s+', text.strip()) if s.strip()]

def clean_text(text: str) -> str:
    if not text or not isinstance(text, str): return ""
    text = re.sub(r'```math\n\d+```|```math\n\d+,\s*\d+```|```math\n\d+–\d+```', '', text)
    text = re.sub(r'\b[A-Z][a-z]+,\s*[A-Z]\.\s*', '', text)
    text = re.sub(r'\b[A-Z]\.\.\s*', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def create_download_file(content: str, file_type: str = "txt") -> bytes:
    """Create a downloadable file in TXT or PDF format."""
    try:
        if file_type == "pdf":
            if not importlib.util.find_spec("reportlab"):
                return content.encode('utf-8')
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            flowables = [Paragraph(line.replace('\n', '<br/>'), styles['Normal']) for line in content.split('\n') if line.strip()]
            doc.build(flowables)
            buffer.seek(0)
            return buffer.getvalue()
        return content.encode('utf-8')
    except Exception:
        return content.encode('utf-8')

def auto_select_summary_mode(text: str) -> str:
    """Determine summary mode based on text length."""
    wc = len(re.findall(r"\b[\w'-]+\b", text))
    return 'extractive' if wc <= 400 else 'hybrid'

@st.cache_data
def build_qa_context(text: str, question: str, chunk_size: int = 5, max_chars: int = 4000) -> str:
    """Build context for Q&A by selecting relevant sentences."""
    try:
        model = get_st_model()
        if not model:
            return clean_text(text)[:max_chars]
        sentences = [s for s in split_sentences(text) if len(s) > 10][:1500]
        if not sentences:
            return clean_text(text)[:max_chars]
        q_emb = model.encode(question, convert_to_tensor=True, normalize_embeddings=True)
        s_emb = model.encode(sentences, convert_to_tensor=True, normalize_embeddings=True)
        sims = util.cos_sim(q_emb, s_emb)[0].tolist()
        ranked = [s for _, s in sorted(zip(sims, sentences), key=lambda x: x[0], reverse=True)]
        ctx = ""
        for s in ranked[:chunk_size * 4]:
            if len(ctx) + len(s) + 1 > max_chars:
                break
            ctx += ("\n" if ctx else "") + s
        return clean_text(ctx) if ctx else clean_text(text)[:max_chars]
    except Exception:
        return clean_text(text)[:max_chars]

def is_bad_answer(answer: str, question: str) -> bool:
    """Check if an answer is invalid or too similar to the question."""
    a = (answer or "").strip().lower()
    q = (question or "").strip().lower()
    if not a or len(a) < 3 or a == q or SequenceMatcher(None, a, q).ratio() > 0.85:
        return True
    return False

def pick_best_result(results: List[dict], question: str, desired_sentences: int = 3) -> dict:
    """Select the best result based on confidence and sentence count."""
    cleaned = [r for r in results if r and r.get('answer') and not is_bad_answer(r['answer'], question)]
    if cleaned:
        for r in cleaned:
            sentence_count = len(split_sentences(r['answer']))
            length_bonus = min(20, len(r['answer']) / 50) if question else 0
            sentence_diff = abs(sentence_count - desired_sentences) if not question else 0
            cohere_bonus = 20 if not question and r.get('method', '').startswith('cohere') else 0
            r['adjusted_confidence'] = r.get('confidence', 0) + length_bonus - (sentence_diff * 10) + cohere_bonus
        best = max(cleaned, key=lambda x: x.get('adjusted_confidence', 0))
        return best
    return {'answer': 'Unable to generate answer.', 'confidence': 0, 'method': 'none'}

# -----------------------------
# Provider Wrappers
# -----------------------------
@st.cache_data(show_spinner=False)
def cached_query_hf_local(question: str, context: str, sentences: int = 3):
    """Query local Hugging Face model for summarization or Q&A."""
    context = clean_text(context or "").strip()
    if not context:
        return {"answer": "", "confidence": 0, "method": "hf_local_none"}
    try:
        if not hf_tokenizer:
            return None
        tokens = hf_tokenizer(context, truncation=True, max_length=1000, return_tensors="pt")
        truncated_context = hf_tokenizer.decode(tokens["input_ids"][0], skip_special_tokens=True)
        if not question:
            summarizer = get_hf_summarizer()
            if not summarizer:
                return None
            max_length = max(50, sentences * 20)
            min_length = max(20, sentences * 10)
            result = summarizer(truncated_context, max_length=max_length, min_length=min_length, do_sample=False)
            return {'answer': clean_text(result[0].get('summary_text', '')).strip(), 'confidence': 85, 'method': 'hf_local_summarization'}
        qa_model = get_hf_qa()
        if not qa_model:
            return None
        result = qa_model(question=question, context=truncated_context)
        return {'answer': clean_text(result.get('answer', '')).strip(), 'confidence': result.get('score', 0) * 100, 'method': 'hf_local_qa'}
    except Exception:
        return None

@st.cache_data(show_spinner=False)
def query_hf_api(question: str, context: str):
    """Query Hugging Face API for summarization or Q&A."""
    if not HF_API_KEY:
        return None
    context = clean_text(context or "").strip()
    if not context:
        return None
    try:
        if not hf_tokenizer:
            return None
        tokens = hf_tokenizer(context, truncation=True, max_length=1000, return_tensors="pt")
        truncated_context = hf_tokenizer.decode(tokens["input_ids"][0], skip_special_tokens=True)
        headers = {"Authorization": f"Bearer {HF_API_KEY}"}
        if question:
            url = "https://api-inference.huggingface.co/models/deepset/roberta-base-squad2"
            payload = {"inputs": {"question": question, "context": truncated_context}}
        else:
            url = "https://api-inference.huggingface.co/models/facebook/bart-large-cnn"
            payload = {"inputs": truncated_context, "parameters": {"max_length": 150, "min_length": 30}}
        r = requests.post(url, headers=headers, json=payload, timeout=60)
        r.raise_for_status()
        result = r.json()
        if question:
            if isinstance(result, dict) and 'answer' in result:
                return {'answer': clean_text(result['answer']).strip(), 'confidence': result.get('score', 0.9) * 100, 'method': 'hf_api_qa'}
            if isinstance(result, list) and result:
                return {'answer': clean_text(result[0].get('answer', '')).strip(), 'confidence': result[0].get('score', 0.9) * 100, 'method': 'hf_api_qa'}
        else:
            if isinstance(result, list) and result:
                return {'answer': clean_text(result[0].get('summary_text', '')).strip(), 'confidence': 90, 'method': 'hf_api_sum'}
            if isinstance(result, dict) and 'summary_text' in result:
                return {'answer': clean_text(result.get('summary_text', '')).strip(), 'confidence': 90, 'method': 'hf_api_sum'}
        return None
    except Exception:
        return None

@st.cache_data(show_spinner=False)
def cached_query_cohere(question: str, context: str, sentences: int = 3, _cache_key: str = ""):
    """Query Cohere API for summarization or Q&A with unique cache key."""
    if not COHERE_API_KEY:
        return None
    context = clean_text(context or "").strip()
    if not context or len(context) < 10:
        return None
    try:
        co = cohere.Client(COHERE_API_KEY)
        if not hf_tokenizer:
            return None
        tokens = hf_tokenizer(context, truncation=True, max_length=4000, return_tensors="pt")
        truncated_context = hf_tokenizer.decode(tokens["input_ids"][0], skip_special_tokens=True)
        if len(truncated_context.strip()) < 10:
            return None
        for attempt in range(5):
            try:
                if question:
                    response = co.chat(
                        message=f"Based on the following text, provide a concise answer to the user's question.\n\nText:\n\"\"\"\n{truncated_context}\n\"\"\"\n\nQuestion: {question}",
                        model="command-r-plus-08-2024",
                        temperature=0.2,
                        max_tokens=500
                    )
                    answer = clean_text(response.text).strip()
                    if not answer:
                        return None
                    return {'answer': answer, 'confidence': 90, 'method': 'cohere_qa'}
                response = co.chat(
                    message=f"Summarize the following text in exactly {sentences} sentences.\n\nText:\n\"\"\"\n{truncated_context}\n\"\"\"",
                    model="command-r-plus-08-2024",
                    temperature=0.3,
                    max_tokens=1000
                )
                answer = clean_text(response.text).strip()
                if not answer:
                    return None
                sentences_list = split_sentences(answer)
                if len(sentences_list) > sentences:
                    answer = ' '.join(sentences_list[:sentences])
                    if not answer.endswith('.'):
                        answer += '.'
                return {'answer': answer, 'confidence': 90, 'method': 'cohere_summarization'}
            except cohere.errors.CohereAPIError as e:
                if e.status_code == 429 and attempt < 4:
                    time.sleep(2 ** attempt)
                    continue
                return None
    except Exception:
        return None

# -----------------------------
# Document Processing
# -----------------------------
class DocProcessor:
    """Handle document extraction for various file types."""
    def extract_text_from_file(self, uploaded_file, force_ocr: bool = False, default_max_pages: int = 30) -> str:
        if not uploaded_file:
            return "No file uploaded."
        file_name = uploaded_file.name.lower()
        if len(uploaded_file.getvalue()) > 10 * 1024 * 1024:
            return "File too large (max 10MB)."
        if file_name.endswith('.pdf'):
            try:
                pdf_reader = PdfReader(uploaded_file)
                text = ""
                for page in pdf_reader.pages[:default_max_pages]:
                    page_text = page.extract_text() or ""
                    text += page_text + "\n"
                text = clean_text(text)
                if len(pdf_reader.pages) > default_max_pages:
                    st.caption(f"Processed first {default_max_pages} PDF pages (auto limit).")
                return text or "No text extracted from PDF."
            except Exception:
                return "No text extracted."
        if file_name.endswith('.txt'):
            try:
                return clean_text(uploaded_file.getvalue().decode('utf-8', errors='ignore').strip())
            except Exception:
                return "Text extraction error."
        if file_name.endswith(('.doc', '.docx')):
            try:
                doc = Document(uploaded_file)
                text = '\n'.join(p.text for p in doc.paragraphs)
                return clean_text(text)
            except Exception:
                return "DOC/DOCX extraction error."
        if file_name.endswith(('.ppt', '.pptx')):
            try:
                prs = Presentation(uploaded_file)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
                            parts.append(shape.text_frame.text)
                text = '\n'.join(parts)
                return clean_text(text)
            except Exception:
                return "PPT/PPTX extraction error."
        return "Unsupported file type."

    def get_advanced_stats(self, text: str):
        """Compute advanced text statistics."""
        if not text or not isinstance(text, str):
            return {
                'words': 0, 'sentences': 0, 'paragraphs': 0,
                'characters': 0, 'complexity': 0.0, 'unique_words': 0,
                'reading_time': 0
            }
        words = re.findall(r"\b[\w'-]+\b", text)
        sentences = split_sentences(text)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        long_words = [w for w in words if len(w) > 6]
        complexity = len(long_words) / max(len(words), 1) * 100
        return {
            'words': len(words),
            'sentences': len(sentences),
            'paragraphs': len(paragraphs),
            'characters': len(text),
            'complexity': round(complexity, 1),
            'unique_words': len(set(w.lower() for w in words)),
            'reading_time': max(1, len(words) // 200)
        }

# -----------------------------
# Summarization Engine
# -----------------------------
class QuantumSummarizer:
    """Summarization engine with extractive, abstractive-like, and hybrid modes."""
    def __init__(self):
        self.summarizer = LsaSummarizer()

    def sumy_extractive_summary(self, text: str, sentences: int = 3) -> dict:
        """Generate extractive summary using Sumy LSA."""
        try:
            parser = PlaintextParser.from_string(text, SumyTokenizer("english"))
            summary_sentences = self.summarizer(parser.document, sentences)
            summary = '. '.join(str(s) for s in summary_sentences)
            if summary and not summary.endswith('.'):
                summary += '.'
            return {'summary': clean_text(summary) or 'No summary.', 'confidence': 80, 'method': 'sumy_lsa'}
        except Exception:
            return {'summary': 'Summarization failed.', 'confidence': 0, 'method': 'sumy_error'}

    def abstractive_like(self, text: str, sentences: int = 3) -> dict:
        """Generate abstractive-like summary by selecting and modifying key sentences."""
        raw_sentences = [s.strip() for s in split_sentences(text) if len(s.strip()) > 15]
        keywords = self._extract_key_concepts(text)
        scored = [(self._concept_score(s, keywords), s) for s in raw_sentences]
        top = sorted(scored, reverse=True)[:max(2, sentences // 2)]
        summary = '. '.join([self._abstract_sentence(s, keywords) for _, s in top])
        if summary and not summary.endswith('.'):
            summary += '.'
        conf = min(95, (sum(score for score, _ in top) / max(len(top), 1)) * 100)
        return {'summary': clean_text(summary) or 'No summary.', 'confidence': round(conf, 1), 'method': 'abstractive_light'}

    def hybrid(self, text: str, sentences: int = 6) -> dict:
        """Combine extractive and abstractive-like summaries."""
        ex = self.sumy_extractive_summary(text, max(2, sentences // 2))
        ab = self.abstractive_like(text, max(2, sentences // 2))
        combined = self._dedupe_sentences(f"{ex['summary']} {ab['summary']}", max_sentences=sentences)
        conf = (ex['confidence'] + ab['confidence']) / 2
        return {'summary': clean_text(combined), 'confidence': round(conf, 1), 'method': 'hybrid'}

    def quantum_summarize(self, text: str, sentences: int = 6) -> dict:
        """Main summarization function with auto-selected mode."""
        if not text or not isinstance(text, str):
            return {'summary': 'No text to summarize.', 'confidence': 0, 'method': 'none'}
        text = clean_text(text)
        mode = auto_select_summary_mode(text)
        if mode == 'extractive':
            return self.sumy_extractive_summary(text, sentences)
        return self.hybrid(text, sentences)

    def _extract_key_concepts(self, text: str) -> List[tuple]:
        """Extract key concepts based on word frequency."""
        words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
        stop_words = {'this', 'that', 'with', 'have', 'will', 'from', 'they', 'been', 'were', 'said', 'into', 'your', 'their'}
        word_freq = {w: words.count(w) for w in set(words) if w not in stop_words}
        return sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]

    def _concept_score(self, sentence: str, keywords: List[tuple]) -> float:
        """Score a sentence based on keyword overlap."""
        s_words = set(re.findall(r'\b[a-zA-Z]{4,}\b', sentence.lower()))
        k_words = set(word for word, _ in keywords)
        return len(s_words.intersection(k_words)) / max(len(s_words), 1)

    def _abstract_sentence(self, sentence: str, keywords: List[tuple]) -> str:
        """Modify a sentence to make it more abstractive-like."""
        words = sentence.split()
        key_concepts = [word for word, _ in keywords[:5]]
        if any(concept in sentence.lower() for concept in key_concepts):
            return ' '.join(words[:min(22, len(words))])
        return sentence

    def _dedupe_sentences(self, text: str, max_sentences: int = 6) -> str:
        """Remove duplicate sentences and limit to max_sentences."""
        sents = split_sentences(text)
        unique = []
        for s in sents:
            if not any(s.lower() in u.lower() or u.lower() in s.lower() for u in unique):
                unique.append(s)
        out = '. '.join(unique[:max_sentences])
        if out and not out.endswith('.'):
            out += '.'
        return out

# -----------------------------
# Heuristic Q&A
# -----------------------------
class NeuroQA:
    """Heuristic-based question answering system."""
    def neural_answer(self, question: str, document: str) -> dict:
        """Generate an answer using heuristic methods."""
        if not question or not document:
            return {'answer': '', 'confidence': 0, 'method': 'neural_error'}
        document = clean_text(document)
        contexts = self._discover_neural_contexts(question, document)
        if not contexts:
            return {'answer': '', 'confidence': 0, 'method': 'neural_no_match'}
        best_context = contexts[0]
        sentences = [s.strip() for s in split_sentences(best_context['text'])]
        question_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', question.lower()))
        scored_sentences = []
        for sentence in sentences:
            sentence_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', sentence.lower()))
            score = len(question_words.intersection(sentence_words))
            scored_sentences.append((score, sentence))
        top_sentences = sorted(scored_sentences, key=lambda x: x[0], reverse=True)[:3]
        answer = '. '.join([s for _, s in top_sentences if s])
        if answer and not answer.endswith('.'):
            answer += '.'
        max_score = sum(score for score, _ in top_sentences) / max(len(top_sentences), 1)
        return {
            'answer': clean_text(answer).strip(),
            'confidence': min(92, max_score / max(len(question_words), 1) * 100),
            'method': 'neural_synthesis'
        }

    def _discover_neural_contexts(self, question: str, document: str) -> List[dict]:
        """Find relevant contexts for the question."""
        sentences = [s.strip() for s in split_sentences(document) if len(s.strip()) > 10]
        question_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', question.lower()))
        contexts = []
        window_size = 3
        for i in range(len(sentences) - window_size + 1):
            context = '. '.join(sentences[i:i + window_size])
            context_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', context.lower()))
            overlap = len(question_words.intersection(context_words))
            if overlap > 0:
                score = overlap / max(len(question_words), 1)
                if score > 0.2:
                    contexts.append({'text': context, 'score': score})
        return sorted(contexts, key=lambda x: x['score'], reverse=True)[:3]

# -----------------------------
# Keyword Extraction
# -----------------------------
def extract_quantum_keywords(text: str, top_k: int = 10) -> List[tuple]:
    """Extract top keywords based on frequency."""
    text = clean_text(text)
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
    stop_words = {
        'this', 'that', 'with', 'have', 'will', 'from', 'they', 'been', 'were', 'said',
        'each', 'which', 'their', 'time', 'about', 'would', 'there', 'could', 'other',
        'after', 'first', 'well', 'also', 'make', 'here', 'where', 'much', 'take', 'than', 'only'
    }
    filtered = [w for w in words if w not in stop_words]
    freq = {w: filtered.count(w) for w in set(filtered)}
    return sorted(freq.items(), key=lambda x: x[1], reverse=True)[:top_k]

# -----------------------------
# Main App
# -----------------------------
def main():
    """Main application function."""
    # Initialize session state
    if 'doc_processor' not in st.session_state:
        st.session_state.doc_processor = DocProcessor()
    if 'quantum_summarizer' not in st.session_state:
        st.session_state.quantum_summarizer = QuantumSummarizer()
    if 'neuro_qa' not in st.session_state:
        st.session_state.neuro_qa = NeuroQA()
    if 'active_page' not in st.session_state:
        st.session_state.active_page = 'upload'
    if 'document_text' not in st.session_state:
        st.session_state.document_text = ""
    if 'original_file_bytes' not in st.session_state:
        st.session_state.original_file_bytes = None
    if 'original_file_name' not in st.session_state:
        st.session_state.original_file_name = None
    if 'original_mime' not in st.session_state:
        st.session_state.original_mime = None
    if 'neural_history' not in st.session_state:
        st.session_state.neural_history = []
    if 'last_summary' not in st.session_state:
        st.session_state.last_summary = None

    load_revolutionary_css()
    st.markdown('<h1 class="main-title">DOCUVERSE AI</h1>', unsafe_allow_html=True)
    st.markdown('<p class="subtitle">Revolutionary Document Intelligence Platform</p>', unsafe_allow_html=True)

    # Navigation
    top_cols = st.columns(5)
    pages = ['upload', 'text', 'analysis', 'summary', 'qa']
    labels = ["Document Upload", "Text Input", "Analysis", "Summary", "Q&A"]
    for col, page, label in zip(top_cols, pages, labels):
        with col:
            if st.button(label, key=f"top_{page}"):
                st.session_state.active_page = page

    # Upload Page
    if st.session_state.active_page == 'upload':
        st.markdown('<div class="content-section"><h2 class="section-title">Document Upload</h2></div>', unsafe_allow_html=True)
        uploaded_file = st.file_uploader(
            "DRAG YOUR DOCUMENT INTO THE FIELD",
            type=['pdf', 'txt', 'doc', 'docx', 'ppt', 'pptx'],
            key="quantum_uploader"
        )

        if uploaded_file:
            st.session_state.original_file_bytes = uploaded_file.getvalue()
            st.session_state.original_file_name = uploaded_file.name
            st.session_state.original_mime = uploaded_file.type or "application/octet-stream"

            file_size = len(st.session_state.original_file_bytes) / 1024 / 1024
            st.markdown(f"""
            <div class="cyber-card">
                <h4 class="cyber-text">File Detected</h4>
                <p><strong>Filename:</strong> {escape(uploaded_file.name)}</p>
                <p><strong>Size:</strong> {file_size:.1f} MB</p>
                <p><strong>Type:</strong> {escape(st.session_state.original_mime)}</p>
                <p><strong>Status:</strong> <span class="cyber-text">Ready for processing</span></p>
            </div>
            """, unsafe_allow_html=True)

            if st.button("Initiate Extraction", key="quantum_extract"):
                with st.spinner("Processing..."):
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    status_text.text("Analyzing document structure...")
                    progress_bar.progress(25); time.sleep(0.2)
                    status_text.text("Extracting text...")
                    progress_bar.progress(55); time.sleep(0.2)
                    status_text.text("Finalizing...")
                    progress_bar.progress(85)
                    text = st.session_state.doc_processor.extract_text_from_file(uploaded_file, force_ocr=False)
                    progress_bar.progress(100); time.sleep(0.1)
                    status_text.text("Extraction complete!")

                    if text and not text.startswith(("Unsupported file type", "Text extraction error", "DOC/DOCX extraction error", "PPT/PPTX extraction error", "File too large")):
                        st.session_state.document_text = text
                        progress_bar.empty(); status_text.empty()
                        st.success("✅ Document extraction successful.")

                        st.markdown("<h4>Full Document (extracted)</h4>", unsafe_allow_html=True)
                        st.markdown(f'<div class="fulltext-container">{escape(st.session_state.document_text)}</div>', unsafe_allow_html=True)

                        dl_col1, dl_col2 = st.columns(2)
                        with dl_col1:
                            st.download_button(
                                label="Download Extracted Text (TXT)",
                                data=st.session_state.document_text.encode('utf-8'),
                                file_name=f"extracted_text_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                mime="text/plain",
                                key="download_extracted_text"
                            )
                        with dl_col2:
                            st.download_button(
                                label="Download Original File",
                                data=st.session_state.original_file_bytes,
                                file_name=st.session_state.original_file_name or "document",
                                mime=st.session_state.original_mime,
                                key="download_original_file"
                            )
                    else:
                        st.error(text or "Extraction failed.")
                        progress_bar.empty(); status_text.empty()

        if st.button("Reset", key="reset_upload"):
            st.session_state.document_text = ""
            st.session_state.original_file_bytes = None
            st.session_state.original_file_name = None
            st.session_state.original_mime = None
            st.cache_resource.clear()
            st.cache_data.clear()
            st.rerun()

    # Text Input Page
    if st.session_state.active_page == 'text':
        st.markdown('<div class="content-section"><h2 class="section-title">Text Input</h2></div>', unsafe_allow_html=True)
        input_text = st.text_area("Enter your text here:", height=400, placeholder="Paste your document text here...", key="text_input_area")
        if input_text:
            words_count = len(re.findall(r"\b[\w'-]+\b", input_text))
            st.caption(f"Word count: {words_count}/5000")
            if words_count > 5000:
                st.error("Input exceeds 5000-word limit. Please shorten your text.")
            elif st.button("Process Text", key="process_text_btn"):
                st.session_state.document_text = clean_text(input_text)
                st.success("✅ Text processed successfully!")
                st.markdown("<h4>Full Document (entered)</h4>", unsafe_allow_html=True)
                st.markdown(f'<div class="fulltext-container">{escape(st.session_state.document_text)}</div>', unsafe_allow_html=True)
                st.download_button(
                    label="Download Input Text (TXT)",
                    data=st.session_state.document_text.encode('utf-8'),
                    file_name=f"input_text_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                    key="download_input_text"
                )
        if st.button("Reset", key="reset_text"):
            st.session_state.document_text = ""
            st.cache_resource.clear()
            st.cache_data.clear()
            st.rerun()

    # Analysis Page
    if st.session_state.active_page == 'analysis':
        if not st.session_state.document_text:
            st.info("🌌 Please upload or input a document first")
        else:
            st.markdown('<div class="content-section"><h2 class="section-title">Document Analysis</h2></div>', unsafe_allow_html=True)
            stats = st.session_state.doc_processor.get_advanced_stats(st.session_state.document_text)
            c1, c2, c3, c4 = st.columns(4)
            with c1:
                st.markdown(f"""<div class="metric-card"><div class="metric-value">{stats['words']:,}</div><div class="metric-label">Words</div></div>""", unsafe_allow_html=True)
            with c2:
                st.markdown(f"""<div class="metric-card"><div class="metric-value">{stats['sentences']:,}</div><div class="metric-label">Sentences</div></div>""", unsafe_allow_html=True)
            with c3:
                st.markdown(f"""<div class="metric-card"><div class="metric-value">{stats['complexity']:.1f}%</div><div class="metric-label">Complexity</div></div>""", unsafe_allow_html=True)
            with c4:
                st.markdown(f"""<div class="metric-card"><div class="metric-value">{stats['reading_time']}</div><div class="metric-label">Second(s) to read</div></div>""", unsafe_allow_html=True)

            st.markdown("### Key Phrases")
            keywords = extract_quantum_keywords(st.session_state.document_text)
            keyword_html = ''.join(f'<span class="keyword-tag">{escape(word)} ({freq})</span>' for word, freq in keywords)
            st.markdown(f'<div style="text-align: center; margin: 0.6rem 0 0.2rem 0;">{keyword_html}</div>', unsafe_allow_html=True)

            if st.button("Reset", key="reset_analysis"):
                st.session_state.document_text = ""
                st.cache_resource.clear()
                st.cache_data.clear()
                st.rerun()

    # Summary Page
    if st.session_state.active_page == 'summary':
        if not st.session_state.document_text.strip():
            st.error("🌌 No valid document text to summarize. Please upload or input a document.")
        else:
            st.markdown('<div class="content-section"><h2 class="section-title">Summary</h2></div>', unsafe_allow_html=True)
            col_params, col_content = st.columns([1, 2])

            with col_params:
                st.markdown('<div class="cyber-card"><h4 class="cyber-text">Settings</h4></div>', unsafe_allow_html=True)
                length = st.slider("Summary length (sentences):", 2, 15, 8, key="summary_length")

            with col_content:
                if st.button("Generate Summary", key="summary_generate"):
                    with st.spinner("Generating summary..."):
                        context = clean_text(st.session_state.document_text)
                        if not context:
                            st.error("No valid text to summarize after cleaning.")
                            return
                        results = []

                        try:
                            qres = st.session_state.quantum_summarizer.quantum_summarize(context, sentences=length)
                            results.append({'answer': qres['summary'], 'confidence': qres['confidence'], 'method': qres['method']})
                        except Exception:
                            results.append({'answer': '', 'confidence': 0, 'method': 'sumy_error'})

                        hf_local = cached_query_hf_local("", context, sentences=length)
                        if hf_local:
                            results.append(hf_local)

                        top_conf = max(r.get('confidence', 0) for r in results if r) if results else 0
                        if top_conf < 80 and HF_API_KEY:
                            api_res = query_hf_api("", context)
                            if api_res:
                                results.append(api_res)
                        if COHERE_API_KEY:
                            cache_key = f"{hash(context)}_{datetime.now().strftime('%Y%m%d%H%M%S')}"
                            cohere_res = cached_query_cohere("", context, sentences=length, _cache_key=cache_key)
                            if cohere_res:
                                results.append(cohere_res)

                        best = pick_best_result(results, "", desired_sentences=length)
                        st.session_state.last_summary = best
                        summary_text = '\n\n'.join(split_sentences(best.get('answer', '')))

                        method = (best.get('method') or '').lower()
                        engine_label = {
                            'hf_api': 'Hugging Face API',
                            'hf_local': 'Local model',
                            'cohere': 'Cohere AI',
                            'sumy': 'Built-in (Concise/Hybrid)',
                            'none': 'Automatic'
                        }.get(method.split('_')[0], 'Automatic')

                        st.markdown(f"""
                        <div class="cyber-card">
                            <h4 class="cyber-text">Best Summary</h4>
                            <div style="background: rgba(0, 255, 127, 0.05); padding: 1.0rem; border-radius: 12px; margin: 0.6rem 0; border-left: 4px solid #00FF7F;">
                                <p style="font-size: 1.05rem; line-height: 1.75; color: #E2E8F0;">{escape(summary_text)}</p>
                            </div>
                            <div style="display: flex; gap: 1rem; flex-wrap: wrap; margin-top: 0.5rem;">
                                <span class="cyber-text">Confidence: {round(best.get('confidence', 0), 1)}%</span>
                                <span class="cyber-text">Engine: {escape(engine_label)}</span>
                                <span class="cyber-text">Sentences: {len(split_sentences(best.get('answer', '')))}</span>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)

                        file_content = f"""DOCUVERSE AI - SUMMARY REPORT
Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S IST")}
Engine: {engine_label}
Mode: Automatic
Confidence: {round(best.get('confidence',0),1)}%
Sentence Count: {len(split_sentences(best.get('answer', '')))}
SUMMARY:
{summary_text}
---
© 2025 DocuVerse AI - Revolutionary Document Intelligence Platform"""

                        d1, d2 = st.columns(2)
                        with d1:
                            st.download_button(
                                label="Download TXT",
                                data=create_download_file(file_content, "txt"),
                                file_name=f"summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                mime="text/plain",
                                key="download_summary_txt"
                            )
                        with d2:
                            if importlib.util.find_spec("reportlab"):
                                st.download_button(
                                    label="Download PDF",
                                    data=create_download_file(file_content, "pdf"),
                                    file_name=f"summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                                    mime="application/pdf",
                                    key="download_summary_pdf"
                                )

            if st.button("Reset", key="reset_summary"):
                st.session_state.last_summary = None
                st.cache_resource.clear()
                st.cache_data.clear()
                st.rerun()

    # Q&A Page
    if st.session_state.active_page == 'qa':
        if not st.session_state.document_text.strip():
            st.error("🌌 No valid document text to query. Please upload or input a document.")
        else:
            st.markdown('<div class="content-section"><h2 class="section-title">Ask your document</h2></div>', unsafe_allow_html=True)
            question = st.text_input("Question:", placeholder="e.g., What are the key risks and action items?", key="qa_question")
            
            if st.button("Get Answer", key="qa_run"):
                if not question.strip():
                    st.warning("Please enter a valid question.")
                else:
                    with st.spinner("Searching and answering..."):
                        context = clean_text(build_qa_context(st.session_state.document_text, question, chunk_size=5, max_chars=4000))
                        if not context:
                            st.error("No valid context available for Q&A.")
                            return
                        results = []

                        try:
                            neuro = st.session_state.neuro_qa.neural_answer(question, context)
                            results.append(neuro)
                        except Exception:
                            results.append({'answer': '', 'confidence': 0, 'method': 'neural_error'})

                        hf_local = cached_query_hf_local(question, context)
                        if hf_local:
                            results.append(hf_local)

                        top_conf = max([r.get('confidence', 0) for r in results if r], default=0)
                        if top_conf < 75 and HF_API_KEY:
                            api_res = query_hf_api(question, context)
                            if api_res:
                                results.append(api_res)
                        if top_conf < 75 and COHERE_API_KEY:
                            cohere_res = cached_query_cohere(question, context)
                            if cohere_res:
                                results.append(cohere_res)

                        best_result = pick_best_result(results, question, desired_sentences=3)
                        if best_result.get('confidence', 0) < 70 or is_bad_answer(best_result.get('answer', ''), question):
                            st.warning("Low confidence answer. Try rephrasing your question.")

                        answer_text = '\n\n'.join(split_sentences(best_result.get('answer', '')))
                        method = (best_result.get('method') or '').lower()
                        engine_label = {
                            'hf_api': 'Hugging Face API',
                            'hf_local': 'Local model',
                            'cohere': 'Cohere AI',
                            'neural': 'Built-in Heuristic',
                            'none': 'Automatic'
                        }.get(method.split('_')[0], 'Automatic')

                        st.session_state.neural_history.append({
                            'question': question,
                            'answer': best_result.get('answer', ''),
                            'confidence': best_result.get('confidence', 0),
                            'method': best_result.get('method', ''),
                            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S IST")
                        })
                        st.session_state.neural_history = st.session_state.neural_history[-50:]

                        st.markdown(f"""
                        <div class="cyber-card">
                            <h4 class="cyber-text">Answer</h4>
                            <div style="background: rgba(78, 205, 196, 0.05); padding: 1.0rem; border-radius: 12px; margin: 0.6rem 0; border-left: 4px solid #4ECDC4;">
                                <p><strong>Q:</strong> {escape(question)}</p>
                                <p><strong>A:</strong> {escape(answer_text)}</p>
                            </div>
                            <div style="margin-top: 0.4rem; display: flex; gap: 1rem; flex-wrap: wrap;">
                                <span class="cyber-text">Confidence: {round(best_result.get('confidence', 0), 1)}%</span>
                                <span class="cyber-text">Engine: {escape(engine_label)}</span>
                                <span class="cyber-text">Time: {datetime.now().strftime("%Y-%m-%d %H:%M:%S IST")}</span>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)

            if st.session_state.neural_history:
                st.markdown("### 🕒 Recent Q&A")
                for i, qa in enumerate(reversed(st.session_state.neural_history[-5:])):
                    answer_text = '\n\n'.join(split_sentences(qa['answer']))
                    engine_label = {
                        'hf_api': 'Hugging Face API',
                        'hf_local': 'Local model',
                        'cohere': 'Cohere AI',
                        'neural': 'Built-in Heuristic',
                        'none': 'Automatic'
                    }.get((qa.get('method') or '').lower().split('_')[0], 'Automatic')
                    with st.expander(f"💭 {escape(qa['question'][:60])}... ({qa['timestamp']})", expanded=(i == 0)):
                        st.markdown(f"""
                        <div class="cyber-card">
                            <p><strong>❓ Question:</strong> {escape(qa['question'])}</p>
                            <p><strong>🤖 Answer:</strong> {escape(answer_text)}</p>
                            <div style="margin-top: 0.4rem;">
                                <span class="cyber-text">Confidence: {round(qa['confidence'], 1)}%</span>
                                <span class="cyber-text" style="margin-left: 1rem;">Engine: {escape(engine_label)}</span>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)

            if st.button("Reset", key="reset_qa"):
                st.session_state.neural_history = []
                st.cache_resource.clear()
                st.cache_data.clear()
                st.rerun()

if __name__ == "__main__":
    main()
